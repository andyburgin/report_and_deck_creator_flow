from pptx import Presentation
from pptx.util import Inches
import os
from PIL import Image

# Slide layout ids
TITLE_LAYOUT = 0
AGENDA_LAYOUT = 4
SECTION_LAYOUT = 6
TEXT_LAYOUT = 14
BULLET_LAYOUT = 15
PICTURE_LAYOUT = 14
CLOSING_LAYOUT = 19

# Image Position/Size
IMAGE_TOP = 1.8
IMAGE_LEFT = 0.91
IMAGE_HEIGHT = 5
IMAGE_WIDTH = 11.55

prs = Presentation("shell.pptx")

def add_title(title_text, subtitle_text, notes_text):
    title_slide_layout = prs.slide_layouts[TITLE_LAYOUT]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = title_text
    subtitle.text = subtitle_text

    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes_text


def add_section(title_text, notes_text):
    title_slide_layout = prs.slide_layouts[SECTION_LAYOUT]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = title_text

    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes_text


def add_agenda(title_text, agenda_text, notes_text):
    add_text_bullet(title_text, agenda_text, notes_text,AGENDA_LAYOUT,13)


def add_bullets(title_text, bullets_text, notes_text):
    add_text_bullet(title_text, bullets_text, notes_text,BULLET_LAYOUT,16)


def add_text(title_text, body_text, notes_text):
    add_text_bullet(title_text, body_text, notes_text,TEXT_LAYOUT,16)


def add_text_bullet(title_text, bullets_text, notes_text, layout, placeholder):
    title_slide_layout = prs.slide_layouts[layout]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    text = slide.placeholders[placeholder]

    title.text = title_text
    text.text = bullets_text

    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes_text


def add_picture(title_text, image_path, notes_text):
    title_slide_layout = prs.slide_layouts[PICTURE_LAYOUT]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title

    title.text = title_text

    #top = Inches(1.7)
    #left = Inches(1.98)
    #width = Inches(9.4)
    #height = Inches(5.4)

    # resize and position image based on top corner IMAGE_LEFT,IMAGE_TOP and IMAGE_WIDTHxIMAGE_HEIGHT
    img = Image.open(image_path)
    img_size_ratio =  img.size[0] / img.size[1]
    if img_size_ratio > (IMAGE_WIDTH/IMAGE_HEIGHT):
        height =  Inches(IMAGE_WIDTH / img_size_ratio)
        width = Inches(IMAGE_WIDTH)
    else:
        height = Inches(IMAGE_HEIGHT)
        width = Inches(IMAGE_HEIGHT * img_size_ratio)

    top = Inches(IMAGE_TOP + IMAGE_HEIGHT/2) - height/2
    left = Inches(IMAGE_LEFT + IMAGE_WIDTH/2) - width/2

    pic = slide.shapes.add_picture(image_path, left, top, width, height)

    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes_text


def add_closing(title_text, subtitle_text, notes_text):
    title_slide_layout = prs.slide_layouts[CLOSING_LAYOUT]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = title_text
    subtitle.text = subtitle_text

    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes_text


def dump_layout(layout):
    slide = prs.slides.add_slide(prs.slide_layouts[layout])
    for shape in slide.placeholders:
        print('%d %s %s' % (shape.placeholder_format.idx, shape.name, shape.placeholder_format.type))


def save(filename):
    prs.save(filename)


#if __name__ == "__main__":
#    prs = Presentation("shell.pptx")
#    add_title("my presentation", "it's really rather good", "Hello and welcome to my talk\n\nToday I'm going to talkto you about\n* stuff\n* more stuff\n* even more stuff")
#    add_section("This is a Section","In this section I'm going to say stuff")
#    add_bullets("Here are some points", "point 1\npoint 2\npoint 3\npoint 4", "These are my things to say\n\nhere's some more words\nyet more words for ya")
#    add_picture("This is a nice picture", "image.png", "What a smashing picture, let me tell you about it")
#    add_closing("Thanks for listening", "You've been a lovely audience", "Thanks for comming to the talk, if you have any questions I'm around for the rest of the event")
#    dump_layout(BULLET_LAYOUT)
#    prs.save('test.pptx')
